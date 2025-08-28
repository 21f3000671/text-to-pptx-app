
from typing import Dict, Any, List, Tuple, Optional
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from io import BytesIO
from .template_assets import extract_template_images, pick_logo_candidate

def _pick_layout(prs: Presentation) -> Tuple[int, int]:
    """Pick reasonable defaults; fall back gracefully if indexes are invalid."""
    content_idx = 1  # common default
    title_only_idx = 5
    for i, layout in enumerate(prs.slide_layouts):
        name = (layout.name or "").lower()
        if "title and content" in name or "content" in name:
            content_idx = i
        if "title only" in name:
            title_only_idx = i
    # clamp to range
    max_idx = len(prs.slide_layouts) - 1
    content_idx = min(max(content_idx, 0), max_idx)
    title_only_idx = min(max(title_only_idx, 0), max_idx)
    return content_idx, title_only_idx

def _set_title(shape, text: str):
    try:
        if hasattr(shape, "text_frame") and shape.text_frame:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = text
            p.alignment = PP_ALIGN.LEFT
    except Exception:
        pass

def _set_bullets_textbox(prs, slide, bullets: List[str]):
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), prs.slide_width - Inches(2.0), prs.slide_height - Inches(3.0))
    tf = tb.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = b
        p.level = 0

def _set_bullets_placeholder(placeholder, bullets: List[str]) -> bool:
    if not getattr(placeholder, "has_text_frame", False):
        return False
    try:
        tf = placeholder.text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = b
            p.level = 0
        return True
    except Exception:
        return False

def _add_logo_if_any(prs, slide, logo_bytes: bytes):
    if not logo_bytes or not isinstance(logo_bytes, (bytes, bytearray)):
        return
    try:
        slide.shapes.add_picture(BytesIO(logo_bytes), prs.slide_width - Inches(1.6), Inches(0.2), width=Inches(1.2))
    except Exception:
        pass

def _safe_add_slide(prs: Presentation, layout_idx: int):
    """Always return a valid slide even if a given layout index is problematic."""
    try:
        layout = prs.slide_layouts[layout_idx]
    except Exception:
        layout = prs.slide_layouts[0]
    # In extremely rare cases, layout collection might misbehave; ensure it's a SlideLayout
    try:
        return prs.slides.add_slide(layout)
    except Exception:
        # last-resort: try layout 0
        return prs.slides.add_slide(prs.slide_layouts[0])

def build_presentation_from_outline(template_path: str, outline: Dict[str, Any], output_path: str):
    prs = Presentation(template_path)
    content_idx, title_only_idx = _pick_layout(prs)

    # images = extract_template_images(prs)
    # logo_bytes = pick_logo_candidate(images)

    # Title slide
    if outline.get("title"):
        slide = _safe_add_slide(prs, 0)
        if getattr(slide.shapes, "title", None):
            _set_title(slide.shapes.title, str(outline["title"])[:180])
        # _add_logo_if_any(slide, logo_bytes)

    for s in outline.get("slides", []):
        slide = _safe_add_slide(prs, content_idx)

        if getattr(slide.shapes, "title", None):
            _set_title(slide.shapes.title, (s.get("title") or "Slide")[:180])

        bullets = [str(b) for b in s.get("bullets", [])][:12]
        body_set = False
        for ph in getattr(slide, "placeholders", []):
            try:
                idx = getattr(getattr(ph, "placeholder_format", None), "idx", None)
                if idx == 0:
                    continue  # skip title placeholder
                if _set_bullets_placeholder(ph, bullets):
                    body_set = True
                    break
            except Exception:
                continue
        if not body_set:
            _set_bullets_textbox(prs, slide, bullets)

        # Speaker notes—some templates may not have notes master; guard it.
        if "notes" in s:
            try:
                notes = s.get("notes") or ""
                notes_slide = slide.notes_slide
                notes_tf = notes_slide.notes_text_frame
                notes_tf.clear()
                notes_tf.text = notes[:2000]
            except Exception:
                pass

        # _add_logo_if_any(slide, logo_bytes)

    prs.save(output_path)
