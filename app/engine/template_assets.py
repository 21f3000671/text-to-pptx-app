
from typing import List
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_template_images(prs: Presentation) -> List[bytes]:
    blobs: List[bytes] = []

    for master in prs.slide_masters:
        for shp in master.shapes:
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE and hasattr(shp, "image"):
                try:
                    blobs.append(shp.image.blob)
                except Exception:
                    pass
        for layout in master.slide_layouts:
            for shp in layout.shapes:
                if shp.shape_type == MSO_SHAPE_TYPE.PICTURE and hasattr(shp, "image"):
                    try:
                        blobs.append(shp.image.blob)
                    except Exception:
                        pass

    for slide in prs.slides[:3]:
        for shp in slide.shapes:
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE and hasattr(shp, "image"):
                try:
                    blobs.append(shp.image.blob)
                except Exception:
                    pass

    uniq = []
    seen = set()
    for b in blobs:
        h = (len(b), hash(b[:200]))
        if h not in seen:
            seen.add(h)
            uniq.append(b)
    return uniq

def pick_logo_candidate(images: List[bytes]) -> bytes:
    for b in images:
        if 5000 < len(b) < 400000:
            return b
    return b""
